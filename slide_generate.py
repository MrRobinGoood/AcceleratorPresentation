from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

prs = Presentation("d47176a9b920bbe5.pptx")
title_slide_layout = prs.slide_layouts[0]


def title_slide(title_text, subtitle_text=''):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    subtitle.text = subtitle_text
    title.text = title_text


def group_shape_sample(grp_shape, img, text, img_left=Inches(1.0), img_top=Inches(2.5), img_width=Inches(2.0),
                       img_height=Inches(2.0)):
    text_top = img_top + img_height
    text_left = img_left
    text_width = img_width
    text_height = Inches(0.1)

    grp_shape_picture = grp_shape.shapes.add_picture(img, img_left, img_top, img_width, img_height)
    grp_shape_text = grp_shape.shapes.add_textbox(text_left, text_top, text_width, text_height)
    grp_shape_text.text = text


def access_between_pers(stp, count,  shapes, lst_img, lst_text, img_width=Inches(2.0), img_height=Inches(2.5)):
    step = 0
    i = 0
    while i < count:
        grp_shape = shapes.add_group_shape()
        group_shape_sample(grp_shape, lst_img[i], lst_text[i], img_width=img_width, img_height=img_height)
        grp_shape.left = Inches(1 + step)
        step += stp
        i += 1


def contact_info(lst_img=[], lst_text=[]):
    count = len(lst_img)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = 'Контакты'
    shapes = slide.shapes
    if count == 5:
        access_between_pers(2.3, count, shapes, lst_img, lst_text)
    elif count == 4:
        access_between_pers(3.0, count, shapes, lst_img, lst_text, Inches(2.3), Inches(2.3))
    elif count == 3:
        access_between_pers(3.5, count, shapes, lst_img, lst_text, Inches(2.8), Inches(2.8))
    elif count == 2:
        grp_shape = shapes.add_group_shape()
        group_shape_sample(grp_shape, lst_img[0], lst_text[0], img_width=Inches(3.5), img_height=Inches(3.5))
        grp_shape.left = Inches(3.0)
        grp_shape = shapes.add_group_shape()
        group_shape_sample(grp_shape, lst_img[1], lst_text[1], img_width=Inches(3.5), img_height=Inches(3.5))
        grp_shape.left = Inches(7.0)
    elif count == 1:
        grp_shape = shapes.add_group_shape()
        group_shape_sample(grp_shape, lst_img[0], lst_text[0], img_width=Inches(3.5), img_height=Inches(3.5))
        grp_shape.left = Inches(5.0)


def problems_slide(text):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = 'Проблематика'
    tf = slide.placeholders[1].text_frame
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(18)



title_slide('Помоги стартапу сделать свой Pitch Deck', 'ООО «Акселератор Возможностей» при ИНТЦ МГУ «Воробьевы горы»')
contact_info(['photo_2023-07-12_08-24-17.jpg', 'photo_2023-07-12_08-24-17.jpg', 'photo_2023-07-12_08-24-17.jpg'],
             ['ML / Backend-dev\n@alexsnegur\n+7-960-831-63-91', 'ML / Backend-dev\n@alexsnegur\n+7-960-831-63-91', 'ML / Backend-dev\n@alexsnegur\n+7-960-831-63-91'])
problems_slide('Раскроем небольшую тайну венчура - для привлечения денежных средств и защиты своего проекта,стартапу нужен Pitch-Deck.\nЧто такое Pitch-Deck?\nPitch-Deck представляет собой презентацию-тизер проекта/компании для инвесторов, партнеров,\nжурналистов и других заинтересованных лиц. Цель презентации - привлечение дополнительного\nфинансирования (инвестиций) - примеры можно найти по ссылке.\nПочему это проблема?\nПроблема #1. Недостаток средств\nДля многих стартапов ограниченные финансы создают преграду при разработке качественного Pitch Deck.')

prs.save('test.pptx')
